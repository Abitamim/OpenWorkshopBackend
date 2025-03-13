from fastapi import FastAPI, Form, Request, status, File, UploadFile
from fastapi.responses import HTMLResponse, FileResponse, RedirectResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
import pandas as pd
import uvicorn
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.line import LineFormat
from math import ceil

app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

def filter_dates(date_str, start_date_filter, end_date_filter, duration):
    dates = [d.strip() for d in date_str.split(',')]
    valid_dates = []
    for d_str in dates:
        try:
            d = pd.to_datetime(d_str, format='%b-%d-%Y', errors='coerce')
            if pd.isna(d):
                continue
            workshop_end = d + pd.to_timedelta(duration - 1, unit='d')
            if d >= start_date_filter and (end_date_filter is None or workshop_end <= end_date_filter):
                valid_dates.append(d.strftime('%b-%d-%Y'))
        except Exception:
            continue
    return ', '.join(valid_dates)

@app.post("/preview")
async def preview_workshops(
    request: Request,
    start_date: str = Form(...),
    end_date: str = Form(None),
    countries: list[str] = Form(..., alias="countries[]")
):
    start_date_filter = pd.to_datetime(start_date)
    end_date_filter = pd.to_datetime(end_date) if end_date else None

    preview_data = {}
    try:
        for country in countries:
            df = pd.read_excel('workshops_data/clean_workshops.xlsx', sheet_name=country)
            df['Dates Available'] = df.apply(
                lambda row: filter_dates(row['Dates Available'], 
                start_date_filter, 
                end_date_filter, 
                row['Duration (Days)']), 
                axis=1
            )
            df = df[df['Dates Available'] != '']
            if not df.empty:
                preview_data[country] = df.to_dict('records')

        return {"success": True, "data": preview_data}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    print('Request for index page received')
    try:
        with open('workshops_data/available_countries.txt', 'r', encoding='utf-8') as f:
            countries = f.read().splitlines()
    except FileNotFoundError:
        countries = []
    
    return templates.TemplateResponse('index.html', {
        "request": request, 
        "countries": countries
    })

@app.post("/")
async def search_workshops(
    request: Request,
    start_date: str = Form(...),
    end_date: str = Form(None),
    countries: list[str] = Form(..., alias="countries[]"),
    file_type: str = Form(...)  # Changed from file_types to file_type
):
    start_date_filter = pd.to_datetime(start_date)
    end_date_filter = pd.to_datetime(end_date) if end_date else None

    try:
        if file_type == 'excel':
            # Create Excel file
            output = BytesIO()
            print("before writer")
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                for country in countries: 
                    print("country " + country)
                    df = pd.read_excel('workshops_data/clean_workshops.xlsx', sheet_name=country)
                    print("after read")
                    df['Dates Available'] = df.apply(lambda row: filter_dates(row['Dates Available'], start_date_filter, end_date_filter, row['Duration (Days)']), axis=1)
                    # Drop rows that end up with no valid dates.
                    print("after filter")
                    df = df[df['Dates Available'] != '']
                    print("excel df after filter " + str(df))
                    if df.empty:
                        continue  # Skip country if no workshops pass the filter
                    df.to_excel(writer, sheet_name=country, index=False)
            
            output.seek(0)
            filename = f"workshops_by_country.xlsx"
            
            return StreamingResponse(
                output,
                media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                headers={
                    'Content-Disposition': f'attachment; filename="{filename}"',
                    'Access-Control-Expose-Headers': 'Content-Disposition'
                }
            )
            
        elif file_type == 'pptx':
            # Create PowerPoint file
            template_path = "pptx_templates/template.pptx"
            prs = Presentation(template_path)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            rows_per_slide = 17  # maximum data rows per slide

            for country in countries:
                df = pd.read_excel('workshops_data/clean_workshops.xlsx', sheet_name=country)
                df['Dates Available'] = df.apply(lambda row: filter_dates(row['Dates Available'], start_date_filter, end_date_filter, row['Duration (Days)']), axis=1)
                df = df[df['Dates Available'] != '']
                if df.empty:
                    print("skipping")
                    continue  # Skip country if no workshops pass the filter
                print("after continue")
                num_slides = ceil(len(df) / rows_per_slide)

                for slide_num in range(num_slides):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])

                    # Add title (placed at the top)
                    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.0), Inches(12), Inches(0.75))
                    title_text = title.text_frame.add_paragraph()
                    title_text.text = f"Open Workshops in {country}"
                    title_text.font.size = Pt(24)
                    title_text.font.bold = True

                    # Create table with fixed total_rows regardless of data count
                    # Slice data for this slide and add it to the table starting at row index 1.
                    slide_data = df.iloc[slide_num * rows_per_slide : (slide_num + 1) * rows_per_slide]
                    data_rows = len(slide_data)
                    total_rows_slide = data_rows + 1  # header row + data rows
                    fixed_row_height = Inches(0.3)
                    table_height = total_rows_slide * fixed_row_height

                    # Create table with exactly as many rows as needed
                    table_shape = slide.shapes.add_table(total_rows_slide, 3, Inches(0.5), Inches(1), Inches(12), table_height)
                    table = table_shape.table
                    table.columns[0].width = Inches(6)
                    table.columns[1].width = Inches(2)
                    table.columns[2].width = Inches(4)

                    # Set each row's height
                    for r in range(total_rows_slide):
                        table.rows[r].height = fixed_row_height

                    # Add header row
                    headers = ["Workshop Title", "Duration (Days)", "Dates Available"]
                    for i, header in enumerate(headers):
                        cell = table.cell(0, i)
                        cell.text = header
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(12)
                        paragraph.alignment = PP_ALIGN.CENTER

                    # Fill table with data rows
                    for idx, (_, row_data) in enumerate(slide_data.iterrows(), start=1):
                        for col_idx, value in enumerate(row_data):
                            cell = table.cell(idx, col_idx)
                            cell.text = str(value)
                            paragraph = cell.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(10)
                            if col_idx == 1:  # center-align Duration column
                                paragraph.alignment = PP_ALIGN.CENTER

            ppt_buffer = BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            filename = f"workshops_by_country.pptx"
                    
            return StreamingResponse(
                        ppt_buffer,
                        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        headers={
                            'Content-Disposition': f'attachment; filename="{filename}"',
                            'Access-Control-Expose-Headers': 'Content-Disposition'
                        }
                    )

    except Exception as e:
        print(f"Error: {str(e)}")
        return templates.TemplateResponse("index.html", {
            "request": request,
            "message": f"Error: {str(e)}",
            "success": False,
            "countries": countries
        })

@app.get('/admin', response_class=HTMLResponse)
async def admin(request: Request):
    print('Request for admin page recieved')
    return templates.TemplateResponse('admin.html', {"request": request})

@app.post('/admin', response_class=HTMLResponse)
async def handle_upload(request: Request, file: UploadFile = File(...)):
    try:
        # Read the Excel file
        contents = await file.read()
        
        # You can process the Excel file here
        df = pd.read_excel(BytesIO(contents))
        print("after read")
        df.columns = df.columns.str.strip()
        print('1')
        df = df[df["SubStatus"] == "Open"]
        df['Start Date'] = pd.to_datetime(df['StartDate'], format='%Y-%m-%d')
        df['End Date'] = pd.to_datetime(df['EndDate'], format='%Y-%m-%d')
        print('2')
        df['Workshop Title'] = df['WorkshopTitle']
        print('3')
        df['Duration (Days)'] = (df['End Date'] - df['Start Date']).dt.days + 1
        df['Dates Available'] = df['Start Date'].dt.strftime('%Y-%m-%d')
        df['Delivery Method'] = df['DeliveryMethod']
        df['Delivery Language'] = df['DeliveryLanguage']
        df['Time Zone'] = df['TimeZone']
        countries = df['Country'].unique().tolist()
        with open('workshops_data/available_countries.txt', 'w', encoding='utf-8') as f:
            f.write('\n'.join(countries))
        country_grouped_dfs = {}

        def sort_dates(dates_str):
            # Split the string by commas
            dates_list = dates_str.split(',')
            
            # Convert each date string to a datetime object
            dates_list = [datetime.strptime(date.strip(), '%Y-%m-%d') for date in dates_list]
            
            # Sort the dates in ascending order (closest first)
            dates_list.sort()
        
            # Convert back to string format dd-mmm-yy and join by commas
            return ', '.join([datetime.strftime(date, '%b-%d-%Y') for date in dates_list])
        
        with pd.ExcelWriter("workshops_data/clean_workshops.xlsx") as writer:
            for country in countries:
                country_grouped_dfs[country] = df[df['Country'] == country].copy().groupby('Workshop Title').agg({
                'Duration (Days)': 'first',  # We can just take the first occurrence, assuming length is the same for each group
                'Dates Available': lambda x: ', '.join(sorted(x))  # Combine start dates into a comma-separated list
                }).reset_index()
                country_grouped_dfs[country]['Dates Available'] = country_grouped_dfs[country]['Dates Available'].apply(sort_dates)
                country_grouped_dfs[country].to_excel(writer, index=False, sheet_name=country)
        
        return templates.TemplateResponse('admin.html', {
            "request": request,
            "message": f"Successfully processed {file.filename}",
            "success": True
        })
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        return templates.TemplateResponse('admin.html', {
            "request": request,
            "message": f"Error processing file: {str(e)}",
            "success": False
        })

if __name__ == '__main__':
    uvicorn.run('myapp:app', host='0.0.0.0', port=8000)

